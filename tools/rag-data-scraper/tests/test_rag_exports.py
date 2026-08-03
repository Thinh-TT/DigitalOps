import csv
from datetime import datetime, timezone
import hashlib
import io
import json
from pathlib import Path
import sqlite3
from types import SimpleNamespace
from uuid import uuid4
from zipfile import ZipFile
from xml.etree import ElementTree

from fastapi.testclient import TestClient
import pytest
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from rag_data_scraper.exporters.rag_exporter import (
    ExportTooLargeError,
    ExportDependencyUnavailableError,
    InvalidStagingPackageError,
    RagExportFormat,
    RagExportService,
)
from rag_data_scraper.exporters.staging_exporter import StagingExporter
from rag_data_scraper.db.state_store import CrawlerStateStore
from rag_data_scraper.models.chunk import Chunk, ChunkACL, ChunkSet
from rag_data_scraper.models.observation import (
    DocumentIdentityStrategy,
    DocumentObservation,
    ExtractionQuality,
    QualityStatus,
)
from rag_data_scraper.web import app as web_app


def _digest(value: bytes) -> str:
    return hashlib.sha256(value).hexdigest()


def _create_staging_package(
    tmp_path: Path,
    text: str = "=SUM(1,1)\nNội dung RAG tiếng Việt: pháp luật, dữ liệu.",
) -> tuple[Path, DocumentObservation]:
    job_id = "JOB_EXPORT"
    source_dir = tmp_path / "source"
    source_dir.mkdir(parents=True)
    raw_bytes = b"<html><body>RAG export fixture</body></html>"
    raw_path = source_dir / "document.html"
    normalized_path = source_dir / "normalized.txt"
    raw_path.write_bytes(raw_bytes)
    normalized_path.write_bytes(text.encode("utf-8"))

    observation = DocumentObservation(
        job_id=job_id,
        source_id="fixture",
        source_namespace="example.gov.vn",
        authority_namespace="gov.vn",
        document_identity_strategy=DocumentIdentityStrategy.CONTENT_ONLY,
        canonical_document_key="fixture:document-1",
        source_document_url="https://example.gov.vn/document-1",
        title="Tài liệu thử nghiệm",
        raw_artifact_uri=str(raw_path),
        raw_sha256=_digest(raw_bytes),
        mime_type="text/html",
        normalized_text_uri=str(normalized_path),
        normalized_text_sha256=_digest(text.encode("utf-8")),
        char_count=len(text),
        word_count=len(text.split()),
        extraction_quality=ExtractionQuality(
            status=QualityStatus.CLEAN,
            ocr_used=False,
            confidence_score=1.0,
        ),
        document_metadata={"department": "DigitalOps"},
        crawled_at=datetime(2026, 8, 3, tzinfo=timezone.utc),
    )
    chunk_set = ChunkSet(
        observation_id=observation.observation_id,
        job_id=job_id,
        chunking_strategy="fixture",
        chunker_version="1.0",
        tokenizer_name="fixture:word-count",
        target_tokens=512,
        overlap_tokens=64,
        total_chunks=1,
        created_at=datetime(2026, 8, 3, tzinfo=timezone.utc),
    )
    chunk = Chunk(
        chunk_set_id=chunk_set.chunk_set_id,
        chunk_index=0,
        text=text,
        token_count=len(text.split()),
        character_start=0,
        character_end=len(text),
        content_sha256=_digest(text.encode("utf-8")),
        heading_path="Mục thử nghiệm",
        page_numbers=[1],
        chunk_acl=ChunkACL(
            allowed_roles=["public"],
            denied_roles=[],
            security_classification="public",
        ),
    )
    now = datetime(2026, 8, 3, tzinfo=timezone.utc)
    job_dir = StagingExporter(tmp_path / "staging").export(
        job_id,
        now,
        now,
        [observation],
        [(chunk_set, [chunk])],
        [],
    )
    return job_dir, observation


def test_builds_all_rag_export_formats_with_stable_contract(
    tmp_path: Path,
) -> None:
    text = "=SUM(1,1)\nNội dung RAG tiếng Việt: pháp luật, dữ liệu."
    job_dir, observation = _create_staging_package(tmp_path, text)
    service = RagExportService(job_dir)
    artifacts = {
        export_format: service.build(export_format)
        for export_format in RagExportFormat
    }

    try:
        jsonl_artifact = artifacts[RagExportFormat.CHUNKS_JSONL]
        record = json.loads(jsonl_artifact.path.read_text(encoding="utf-8"))
        assert record["text"] == text
        assert record["metadata"]["format_version"] == "1.0"
        assert record["metadata"]["canonical_document_key"] == (
            "fixture:document-1"
        )
        assert record["metadata"]["allowed_roles"] == ["public"]

        csv_artifact = artifacts[RagExportFormat.CHUNKS_CSV]
        rows = list(
            csv.DictReader(
                io.StringIO(
                    csv_artifact.path.read_text(encoding="utf-8-sig")
                )
            )
        )
        assert len(rows) == 1
        assert rows[0]["text"] == "'" + text
        assert json.loads(rows[0]["allowed_roles_json"]) == ["public"]

        with ZipFile(
            artifacts[RagExportFormat.STAGING_ZIP].path
        ) as archive:
            names = set(archive.namelist())
            assert {
                "manifest.json",
                "document-observations.jsonl",
                "chunk-sets.jsonl",
                "chunks.jsonl",
                "crawler-errors.jsonl",
            }.issubset(names)
            assert "preview.html" not in names
            assert any(
                name.startswith(
                    f"artifacts/{observation.observation_id}/"
                )
                for name in names
            )
            assert all(
                not Path(name).is_absolute() and ".." not in Path(name).parts
                for name in names
            )

        with ZipFile(
            artifacts[RagExportFormat.DOCUMENTS_MARKDOWN_ZIP].path
        ) as archive:
            names = archive.namelist()
            assert "export-manifest.json" in names
            document_name = (
                f"documents/{observation.observation_id}.md"
            )
            markdown = archive.read(document_name).decode("utf-8")
            assert markdown.startswith("---\nformat_version: \"1.0\"")
            assert 'title: "Tài liệu thử nghiệm"' in markdown
            assert markdown.endswith(text + "\n")
        html_text = artifacts[
            RagExportFormat.DOCUMENTS_HTML
        ].path.read_text(encoding="utf-8")
        assert "<title>RAG export JOB_EXPORT</title>" in html_text
        assert "Tài liệu thử nghiệm" in html_text
        assert text in html_text

        json_export = json.loads(
            artifacts[RagExportFormat.CHUNKS_JSON].path.read_text(
                encoding="utf-8"
            )
        )
        assert json_export["format_version"] == "1.0"
        assert json_export["total_chunks"] == 1
        assert json_export["chunks"][0]["text"] == text

        with ZipFile(
            artifacts[RagExportFormat.DOCUMENTS_TXT_ZIP].path
        ) as archive:
            txt_name = (
                f"documents/{observation.observation_id}.txt"
            )
            assert archive.read(txt_name).decode("utf-8") == text
            manifest = json.loads(
                archive.read("export-manifest.json")
            )
            assert manifest["total_documents"] == 1

        pdf = PdfReader(
            artifacts[RagExportFormat.DOCUMENTS_PDF].path
        )
        pdf_text = "\n".join(
            page.extract_text() or "" for page in pdf.pages
        )
        assert "Tài liệu thử nghiệm" in pdf_text
        assert "Nội dung RAG tiếng Việt" in pdf_text

        word = Document(
            artifacts[RagExportFormat.DOCUMENTS_DOCX].path
        )
        word_text = "\n".join(
            paragraph.text for paragraph in word.paragraphs
        )
        assert "Tài liệu thử nghiệm" in word_text
        assert text in word_text

        workbook = load_workbook(
            artifacts[RagExportFormat.CHUNKS_XLSX].path,
            read_only=True,
            data_only=False,
        )
        try:
            chunk_rows = list(
                workbook["Chunks"].iter_rows(values_only=True)
            )
            document_rows = list(
                workbook["Documents"].iter_rows(values_only=True)
            )
            assert chunk_rows[0][0:3] == (
                "id",
                "text",
                "text_part_index",
            )
            assert chunk_rows[1][1] == "'" + text
            assert document_rows[1][3] == "Tài liệu thử nghiệm"
        finally:
            workbook.close()

        deck = Presentation(
            artifacts[RagExportFormat.DOCUMENTS_PPTX].path
        )
        slide_text = "\n".join(
            shape.text
            for slide in deck.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        assert "Tài liệu thử nghiệm" in slide_text
        assert "Nội dung RAG tiếng Việt" in slide_text

        xml_root = ElementTree.parse(
            artifacts[RagExportFormat.DOCUMENTS_XML].path
        ).getroot()
        assert xml_root.attrib["job-id"] == "JOB_EXPORT"
        assert (
            xml_root.findtext("document/normalized-text")
            == text
        )
        assert xml_root.findtext("document/chunks/chunk") == text

        with ZipFile(
            artifacts[RagExportFormat.DOCUMENTS_SVG_ZIP].path
        ) as archive:
            svg_name = (
                f"documents/{observation.observation_id}.svg"
            )
            svg_root = ElementTree.fromstring(
                archive.read(svg_name)
            )
            svg_metadata = svg_root.findtext(
                "{http://www.w3.org/2000/svg}metadata"
            )
            assert svg_metadata is not None
            assert json.loads(svg_metadata)["normalized_text"] == text
    finally:
        for artifact in artifacts.values():
            artifact.cleanup()


def test_rejects_staging_path_traversal_and_tampered_hash(
    tmp_path: Path,
) -> None:
    job_dir, _ = _create_staging_package(tmp_path)
    observations_path = job_dir / "document-observations.jsonl"
    observation = json.loads(observations_path.read_text(encoding="utf-8"))
    observation["raw_artifact_uri"] = "../outside.html"
    observations_path.write_text(
        json.dumps(observation) + "\n",
        encoding="utf-8",
    )

    with pytest.raises(InvalidStagingPackageError):
        RagExportService(job_dir).build(RagExportFormat.CHUNKS_JSONL)


def test_persistent_export_is_checksum_verified_and_not_cleaned_up(
    tmp_path: Path,
) -> None:
    job_dir, _ = _create_staging_package(tmp_path)
    service = RagExportService(job_dir)
    artifact = service.build_persistent(RagExportFormat.CHUNKS_JSONL)

    assert artifact.temporary is False
    assert artifact.path.parent == job_dir / "exports"
    checksum_path = artifact.path.with_name(artifact.path.name + ".sha256")
    assert checksum_path.is_file()
    assert checksum_path.read_text(encoding="ascii").split()[0] == _digest(
        artifact.path.read_bytes()
    )
    assert service.persisted(RagExportFormat.CHUNKS_JSONL) == artifact

    artifact.cleanup()
    assert artifact.path.is_file()
    artifact.path.write_bytes(b"tampered")
    assert service.persisted(RagExportFormat.CHUNKS_JSONL) is None


def test_rejects_export_when_package_exceeds_limit(
    tmp_path: Path,
) -> None:
    job_dir, _ = _create_staging_package(tmp_path)
    with pytest.raises(ExportTooLargeError):
        RagExportService(
            job_dir,
            max_export_bytes=1,
        ).build(RagExportFormat.CHUNKS_JSONL)


def test_export_api_lists_and_downloads_formats(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    job_dir, _ = _create_staging_package(tmp_path)
    monkeypatch.setattr(web_app, "STAGING_DIR", job_dir.parent)
    web_app.JOB_STATUS_MAP.clear()

    with TestClient(web_app.app) as client:
        jobs_response = client.get("/api/jobs")
        assert jobs_response.status_code == 200
        job = jobs_response.json()["jobs"][0]
        assert set(job["export_formats"]) == {
            export_format.value for export_format in RagExportFormat
        }

        formats_response = client.get("/api/jobs/JOB_EXPORT/exports")
        assert formats_response.status_code == 200
        assert len(formats_response.json()["formats"]) == len(RagExportFormat)

        response = client.get(
            "/api/jobs/JOB_EXPORT/exports/chunks_jsonl"
        )
        assert response.status_code == 200
        assert response.headers["cache-control"] == "no-store"
        assert response.headers["x-content-type-options"] == "nosniff"
        assert "JOB_EXPORT-chunks.jsonl" in (
            response.headers["content-disposition"]
        )
        assert json.loads(response.content)["text"].startswith("=SUM")

        unknown_response = client.get(
            "/api/jobs/JOB_EXPORT/exports/not-a-format"
        )
        assert unknown_response.status_code == 422

    web_app.JOB_STATUS_MAP.clear()


def test_export_api_hides_integrity_failure_details(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    job_dir, _ = _create_staging_package(tmp_path)
    chunks_path = job_dir / "chunks.jsonl"
    chunk = json.loads(chunks_path.read_text(encoding="utf-8"))
    chunk["text"] = "tampered"
    chunks_path.write_text(json.dumps(chunk) + "\n", encoding="utf-8")
    monkeypatch.setattr(web_app, "STAGING_DIR", job_dir.parent)

    with TestClient(web_app.app) as client:
        response = client.get(
            "/api/jobs/JOB_EXPORT/exports/chunks_jsonl"
        )

    assert response.status_code == 409
    assert response.json() == {
        "detail": (
            "The staging package is incomplete or failed integrity "
            "validation."
        )
    }


def test_dashboard_exposes_export_controls_and_global_escape_helper() -> None:
    dashboard = (
        Path(web_app.__file__).parent / "static" / "index.html"
    ).read_text(encoding="utf-8")
    assert dashboard.index("function escapeHtml") < dashboard.index(
        "function setPreset"
    )
    assert "function downloadExport" in dashboard
    assert "documents_markdown_zip" in dashboard
    assert "/exports/" in dashboard
    assert 'id="outputFormatSelect"' in dashboard
    assert "export_format: exportFormat" in dashboard
    assert 'id="paginationLimitInput"' in dashboard
    assert "max_pagination_pages: paginationLimit" in dashboard

    assert dashboard.index("async function deleteJob") < dashboard.index(
        "function downloadExport"
    ) < dashboard.index(
        "document.getElementById('crawlForm').addEventListener"
    )
    for export_format in RagExportFormat:
        assert export_format.value in dashboard


def test_export_rejects_chunk_above_declared_hard_limit(
    tmp_path: Path,
) -> None:
    job_dir, _ = _create_staging_package(
        tmp_path,
        "mot hai ba bon",
    )
    chunk_set_path = job_dir / "chunk-sets.jsonl"
    chunk_set = json.loads(chunk_set_path.read_text(encoding="utf-8"))
    chunk_set.update(
        {
            "target_tokens": 1,
            "soft_max_tokens": 1,
            "max_tokens": 1,
            "overlap_tokens": 0,
        }
    )
    chunk_set_path.write_text(
        json.dumps(chunk_set, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    with pytest.raises(InvalidStagingPackageError):
        RagExportService(job_dir).build(RagExportFormat.CHUNKS_JSONL)


def test_create_job_forwards_pre_crawl_export_selection(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    captured: dict[str, object] = {}

    async def fake_execute(**kwargs) -> None:
        captured.update(kwargs)

    monkeypatch.setattr(web_app, "STAGING_DIR", tmp_path / "staging")
    monkeypatch.setattr(web_app, "execute_crawl_job", fake_execute)
    web_app.JOB_STATUS_MAP.clear()

    with TestClient(web_app.app) as client:
        response = client.post(
            "/api/jobs",
            json={
                "job_id": "JOB_FORMAT_PICKER",
                "source": "generic_web",
                "urls": ["https://1.1.1.1/document"],
                "limit": 3,
                "max_pagination_pages": 7,
                "download_attachments": False,
                "export_format": "documents_pdf",
            },
        )

    assert response.status_code == 200
    assert response.json()["export_format"] == "documents_pdf"
    assert captured["export_format"] == RagExportFormat.DOCUMENTS_PDF
    assert captured["download_attachments"] is False
    assert captured["max_pagination_pages"] == 7


def test_delete_job_removes_staging_raw_and_durable_state(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    job_id = "JOB_DELETE"
    staging_base = tmp_path / "staging"
    raw_base = tmp_path / "raw"
    state_db = tmp_path / "state" / "crawler.db"
    staging_job = staging_base / job_id
    raw_job = raw_base / job_id
    staging_job.mkdir(parents=True)
    raw_job.mkdir(parents=True)
    (staging_job / "preview.html").write_text("preview", encoding="utf-8")
    (raw_job / "checkpoint.json").write_text("{}", encoding="utf-8")

    store = CrawlerStateStore(state_db)
    url = "https://example.gov.vn/document"
    store.start_job(
        job_id=job_id,
        source_id="test",
        source_namespace="example.gov.vn",
        authority_namespace="gov.vn",
        identity_strategy="authoritative",
        base_url=url,
    )
    store.prepare_frontier(job_id, [(url, 0)])
    raw_cache_bytes = b"{}"
    store.record_fetch(
        job_id=job_id,
        source_id="test",
        url=url,
        fetch_status="fetched",
        http_status=200,
        content_hash=hashlib.sha256(raw_cache_bytes).hexdigest(),
        etag='"delete-version"',
        raw_artifact_uri=str((raw_job / "checkpoint.json").resolve()),
    )

    class FakeSettings:
        @classmethod
        def load_from_yaml(cls, path):
            return SimpleNamespace(
                storage=SimpleNamespace(raw_base_dir=raw_base)
            )

    monkeypatch.setattr(web_app, "STAGING_DIR", staging_base)
    monkeypatch.setattr(web_app, "STATE_DB", state_db)
    monkeypatch.setattr(web_app, "Settings", FakeSettings)
    web_app.JOB_STATUS_MAP.clear()

    with TestClient(web_app.app) as client:
        response = client.delete(f"/api/jobs/{job_id}")

    assert response.status_code == 200
    assert not staging_job.exists()
    assert not raw_job.exists()
    with sqlite3.connect(state_db) as connection:
        assert connection.execute(
            "SELECT COUNT(*) FROM CrawlJobs WHERE job_id = ?", (job_id,)
        ).fetchone()[0] == 0
        assert connection.execute(
            "SELECT COUNT(*) FROM CrawlFrontier WHERE job_id = ?", (job_id,)
        ).fetchone()[0] == 0
        assert connection.execute(
            "SELECT etag, content_hash, raw_artifact_uri, fetch_status "
            "FROM CrawledResources WHERE url = ?",
            (url,),
        ).fetchone() == (None, None, None, "pending")


def test_html_export_escapes_untrusted_document_content(
    tmp_path: Path,
) -> None:
    source_text = (
        "<script>alert('xss')</script>\n"
        "<img src=x onerror=alert(1)>"
    )
    job_dir, _ = _create_staging_package(tmp_path, source_text)
    artifact = RagExportService(job_dir).build(
        RagExportFormat.DOCUMENTS_HTML
    )
    try:
        exported = artifact.path.read_text(encoding="utf-8")
        assert "<script>" not in exported
        assert "<img src=x" not in exported
        assert "&lt;script&gt;alert(&#x27;xss&#x27;)&lt;/script&gt;" in (
            exported
        )
        assert "&lt;img src=x onerror=alert(1)&gt;" in exported
    finally:
        artifact.cleanup()


def test_export_api_hides_missing_writer_dependency(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    job_dir, _ = _create_staging_package(tmp_path)
    monkeypatch.setattr(web_app, "STAGING_DIR", job_dir.parent)

    def unavailable_writer(
        self: RagExportService,
        export_format: RagExportFormat,
    ) -> None:
        raise ExportDependencyUnavailableError("private dependency detail")

    monkeypatch.setattr(
        web_app.RagExportService,
        "build",
        unavailable_writer,
    )
    with TestClient(web_app.app) as client:
        response = client.get(
            "/api/jobs/JOB_EXPORT/exports/documents_pdf"
        )

    assert response.status_code == 503
    assert response.json() == {
        "detail": "The requested export writer is unavailable."
    }
    assert "private dependency detail" not in response.text
